import uuid
from typing import List, Optional
import PyPDF2
from pptx import Presentation
from io import BytesIO
import os

from app.models.schemas import ParsedDocument, Page
from app.utils.logging import get_logger

logger = get_logger(__name__)


class DocumentParser:
    """Parse PDF and PowerPoint documents"""
    
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    ALLOWED_EXTENSIONS = {'.pdf', '.ppt', '.pptx'}
    
    @staticmethod
    def validate_file(filename: str, file_size: int) -> tuple[bool, Optional[str]]:
        """Validate file format and size"""
        ext = os.path.splitext(filename)[1].lower()
        
        if ext not in DocumentParser.ALLOWED_EXTENSIONS:
            return False, f"Invalid file format. Allowed: {', '.join(DocumentParser.ALLOWED_EXTENSIONS)}"
        
        if file_size > DocumentParser.MAX_FILE_SIZE:
            return False, f"File too large. Maximum size: {DocumentParser.MAX_FILE_SIZE / (1024*1024)}MB"
        
        return True, None
    
    @staticmethod
    async def parse(file_content: bytes, filename: str) -> ParsedDocument:
        """Parse document and extract content"""
        document_id = str(uuid.uuid4())
        ext = os.path.splitext(filename)[1].lower()
        
        try:
            if ext == '.pdf':
                return await DocumentParser._parse_pdf(file_content, document_id, filename)
            elif ext in ['.ppt', '.pptx']:
                return await DocumentParser._parse_pptx(file_content, document_id, filename)
            else:
                raise ValueError(f"Unsupported file format: {ext}")
        except Exception as e:
            logger.error(f"Error parsing document {filename}: {e}")
            raise
    
    @staticmethod
    async def _parse_pdf(content: bytes, document_id: str, filename: str) -> ParsedDocument:
        """Parse PDF file"""
        pdf_file = BytesIO(content)
        reader = PyPDF2.PdfReader(pdf_file)
        
        pages = []
        raw_text = ""
        
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            raw_text += text + "\n"
            
            pages.append(Page(
                page_number=page_num,
                text=text,
                images=[],
                layout=None
            ))
        
        return ParsedDocument(
            document_id=document_id,
            pages=pages,
            metadata={
                "filename": filename,
                "total_pages": len(pages),
                "file_type": "pdf"
            },
            raw_text=raw_text
        )
    
    @staticmethod
    async def _parse_pptx(content: bytes, document_id: str, filename: str) -> ParsedDocument:
        """Parse PowerPoint file"""
        pptx_file = BytesIO(content)
        presentation = Presentation(pptx_file)
        
        pages = []
        raw_text = ""
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            raw_text += text + "\n"
            
            pages.append(Page(
                page_number=slide_num,
                text=text,
                images=[],
                layout=None
            ))
        
        return ParsedDocument(
            document_id=document_id,
            pages=pages,
            metadata={
                "filename": filename,
                "total_pages": len(pages),
                "file_type": "pptx"
            },
            raw_text=raw_text
        )
