"""
Enhanced Document Parser - Phase 1: Document Ingestion & Extraction
Uses PyMuPDF for PDF, python-pptx for PowerPoint, with fallback mechanisms
"""
import fitz  # PyMuPDF
import mimetypes
from typing import List, Dict, Any, Optional
from pptx import Presentation
from io import BytesIO
import uuid

from app.utils.logging import get_logger

logger = get_logger(__name__)


class EnhancedDocumentParser:
    """Production-grade document parser with layout intelligence"""
    
    @staticmethod
    def detect_file_type(filename: str, content: bytes) -> str:
        """Phase 1.1: Detect file type"""
        mime_type, _ = mimetypes.guess_type(filename)
        
        # Verify with magic bytes
        if content[:4] == b'%PDF':
            return 'pdf'
        elif content[:2] == b'PK':  # ZIP-based (PPTX)
            return 'pptx'
        
        return mime_type or 'unknown'
    
    @staticmethod
    async def parse_pdf_pymupdf(content: bytes, filename: str) -> Dict[str, Any]:
        """Phase 1.2: PDF extraction with PyMuPDF (fastest + layout coordinates)"""
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            slides = []
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                blocks = []
                
                # Extract text blocks with coordinates
                text_blocks = page.get_text("dict")["blocks"]
                
                for block in text_blocks:
                    if block.get("type") == 0:  # Text block
                        # Detect if it's a title based on font size
                        lines = block.get("lines", [])
                        if lines:
                            first_span = lines[0].get("spans", [{}])[0]
                            font_size = first_span.get("size", 0)
                            
                            block_type = "title" if font_size > 18 else "body"
                            
                            text = ""
                            for line in lines:
                                for span in line.get("spans", []):
                                    text += span.get("text", "") + " "
                            
                            blocks.append({
                                "type": block_type,
                                "text": text.strip(),
                                "bbox": block.get("bbox", []),
                                "font_size": font_size
                            })
                
                # Extract tables (basic detection)
                tables = page.find_tables()
                for table in tables:
                    blocks.append({
                        "type": "table",
                        "data": table.extract(),
                        "bbox": table.bbox
                    })
                
                slides.append({
                    "slide_no": page_num + 1,
                    "blocks": blocks,
                    "raw_text": page.get_text()
                })
            
            doc.close()
            
            return {
                "document_id": str(uuid.uuid4()),
                "filename": filename,
                "total_slides": len(slides),
                "slides": slides,
                "extraction_method": "pymupdf"
            }
        
        except Exception as e:
            logger.error(f"PyMuPDF extraction failed: {e}")
            raise
    
    @staticmethod
    async def parse_pptx(content: bytes, filename: str) -> Dict[str, Any]:
        """Phase 1.4: PPTX extraction with python-pptx"""
        try:
            prs = Presentation(BytesIO(content))
            slides = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                blocks = []
                
                for shape in slide.shapes:
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text:
                        # Detect type based on position and size
                        block_type = "title" if shape.top < 1000000 else "body"
                        
                        blocks.append({
                            "type": block_type,
                            "text": shape.text,
                            "bbox": [shape.left, shape.top, shape.width, shape.height]
                        })
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = [cell.text for cell in row.cells]
                            table_data.append(row_data)
                        
                        blocks.append({
                            "type": "table",
                            "data": table_data,
                            "bbox": [shape.left, shape.top, shape.width, shape.height]
                        })
                
                # Get raw text
                raw_text = "\n".join([b["text"] for b in blocks if "text" in b])
                
                slides.append({
                    "slide_no": slide_num,
                    "blocks": blocks,
                    "raw_text": raw_text
                })
            
            return {
                "document_id": str(uuid.uuid4()),
                "filename": filename,
                "total_slides": len(slides),
                "slides": slides,
                "extraction_method": "python-pptx"
            }
        
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise
    
    @staticmethod
    async def parse(content: bytes, filename: str) -> Dict[str, Any]:
        """Main parsing entry point with automatic format detection"""
        file_type = EnhancedDocumentParser.detect_file_type(filename, content)
        
        if file_type == 'pdf':
            return await EnhancedDocumentParser.parse_pdf_pymupdf(content, filename)
        elif file_type == 'pptx':
            return await EnhancedDocumentParser.parse_pptx(content, filename)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
